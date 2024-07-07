import asyncio
import os
import json
import re
import openai
from pptx import Presentation
import tkinter as tk
from tkinter import filedialog
from explainer import openaiAgent


def choose_file():
    """
    Opens a file dialog for the user to select a PowerPoint file.

    :return: The path to the selected file.
    """
    root = tk.Tk()
    root.withdraw()

    file_path = filedialog.askopenfilename()
    if file_path:
        return file_path
    else:
        print("No file selected.")
        exit(1)


def clean_text(text):
    """
    Replaces all sequences of one or more whitespace characters in the input string text with a single space.

    :param text: The input text to process.
    :return: Text with single spaces between words.
    """
    return re.sub(r'\s+', ' ', text)


async def read_pptx_file(file_name, key, output_path):
    """
    Reads a PowerPoint (.pptx) file, processes each slide's text, and generates AI responses.

    :param file_name: The name of the PowerPoint file to read.
    :param key: The OpenAI API key for authentication.
    :return: None. Saves the results to a JSON file.
    """
    openai.api_key = key
    prs = Presentation(file_name)
    tasks = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_info = " ".join(clean_text(shape.text) for shape in slide.shapes if hasattr(shape, "text"))

        if slide_info.strip():
            tasks.append((slide_number, openaiAgent.pptx_agent(slide_info)))

    results = []
    responses = await asyncio.gather(*(task[1] for task in tasks))

    for (slide_number, _), solutions in zip(tasks, responses):
        results.append({
            "slide_number": slide_number,
            "solutions": solutions
        })

    save_to_json(results, output_path)



def save_to_json(data, output_path):
    """
    Saves the processed data to output path as JSON file.

    :param data: The data to save.
    :param output_path: output path for saving data
    :return: None.
    """
    json_file = os.path.splitext(output_path)[0] + '.json'
    with open(json_file, 'w') as f:
        json.dump(data, f, indent=4)
